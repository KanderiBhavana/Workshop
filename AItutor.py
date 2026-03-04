import hashlib
import os
from typing import Dict, List, Tuple

import chromadb
import requests
import streamlit as st
from chromadb.config import Settings
from pypdf import PdfReader
from sentence_transformers import SentenceTransformer

# -------- NEW IMPORTS FOR MULTI-FORMAT SUPPORT --------
import pdfplumber
import zipfile
from docx import Document
from pptx import Presentation
import pandas as pd
import json
from bs4 import BeautifulSoup


# ---------------- CONFIG ----------------
DB_DIR = "./chroma_db"
COLLECTION_NAME = "rag_docs"
EMBED_MODEL_NAME = "sentence-transformers/all-MiniLM-L6-v2"
OPENROUTER_API_URL = "https://openrouter.ai/api/v1/chat/completions"


# ---------------- TEXT CHUNKING ----------------
def chunk_text(text: str, chunk_size: int = 700, overlap: int = 120) -> List[str]:
    text = " ".join(text.split())
    if not text:
        return []

    chunks = []
    start = 0
    while start < len(text):
        end = min(start + chunk_size, len(text))
        chunks.append(text[start:end])
        if end == len(text):
            break
        start = end - overlap
    return chunks


# ---------------- FILE READER (UPDATED) ----------------
def read_uploaded_file(uploaded_file) -> str:
    name = uploaded_file.name.lower()

    # -------- PDF (complex + fallback) --------
    if name.endswith(".pdf"):
        text = ""
        try:
            with pdfplumber.open(uploaded_file) as pdf:
                for page in pdf.pages:
                    extracted = page.extract_text()
                    if extracted:
                        text += extracted + "\n"
            if text.strip():
                return text
        except:
            pass

        uploaded_file.seek(0)
        reader = PdfReader(uploaded_file)
        pages = [page.extract_text() or "" for page in reader.pages]
        return "\n".join(pages)

    # -------- DOCX --------
    elif name.endswith(".docx"):
        doc = Document(uploaded_file)
        return "\n".join([para.text for para in doc.paragraphs])

    # -------- TXT / MD --------
    elif name.endswith((".txt", ".md")):
        return uploaded_file.read().decode("utf-8", errors="ignore")

    # -------- CSV --------
    elif name.endswith(".csv"):
        df = pd.read_csv(uploaded_file)
        return df.to_string()

    # -------- Excel --------
    elif name.endswith((".xlsx", ".xls")):
        df = pd.read_excel(uploaded_file)
        return df.to_string()

    # -------- PPTX --------
    elif name.endswith(".pptx"):
        prs = Presentation(uploaded_file)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text

    # -------- JSON --------
    elif name.endswith(".json"):
        data = json.load(uploaded_file)
        return json.dumps(data, indent=2)

    # -------- HTML --------
    elif name.endswith(".html"):
        soup = BeautifulSoup(uploaded_file, "html.parser")
        return soup.get_text()

    # -------- ZIP --------
    elif name.endswith(".zip"):
        text = ""
        with zipfile.ZipFile(uploaded_file, "r") as zip_ref:
            for file_info in zip_ref.infolist():
                filename = file_info.filename.lower()
                if filename.endswith((".txt", ".md", ".json")):
                    with zip_ref.open(file_info) as f:
                        text += f.read().decode("utf-8", errors="ignore") + "\n"
        return text

    # -------- DEFAULT --------
    else:
        return uploaded_file.read().decode("utf-8", errors="ignore")


# ---------------- EMBEDDING ----------------
@st.cache_resource
def get_embedder() -> SentenceTransformer:
    return SentenceTransformer(EMBED_MODEL_NAME)


@st.cache_resource
def get_collection():
    client = chromadb.PersistentClient(
        path=DB_DIR,
        settings=Settings(anonymized_telemetry=False),
    )
    return client.get_or_create_collection(name=COLLECTION_NAME)


# ---------------- STORE DOCUMENTS ----------------
def add_documents(collection, embedder, docs: List[Tuple[str, str]]) -> int:
    all_chunks = []
    ids = []
    metadatas = []

    for source_name, text in docs:
        chunks = chunk_text(text)
        for idx, chunk in enumerate(chunks):
            chunk_id = hashlib.md5(
                f"{source_name}:{idx}:{chunk}".encode("utf-8")
            ).hexdigest()
            ids.append(chunk_id)
            all_chunks.append(chunk)
            metadatas.append({"source": source_name, "chunk": idx})

    if not all_chunks:
        return 0

    existing_ids = set(collection.get(include=[], ids=ids)["ids"])
    new_payload = [
        (i, d, m)
        for i, d, m in zip(ids, all_chunks, metadatas)
        if i not in existing_ids
    ]

    if not new_payload:
        return 0

    new_ids = [item[0] for item in new_payload]
    new_docs = [item[1] for item in new_payload]
    new_meta = [item[2] for item in new_payload]

    embeddings = embedder.encode(
        new_docs,
        normalize_embeddings=True
    ).tolist()

    collection.add(
        ids=new_ids,
        documents=new_docs,
        metadatas=new_meta,
        embeddings=embeddings,
    )

    return len(new_docs)


# ---------------- RETRIEVE ----------------
def retrieve_context(collection, embedder, query: str, top_k: int = 4):
    query_embedding = embedder.encode(
        [query],
        normalize_embeddings=True
    ).tolist()[0]

    results = collection.query(
        query_embeddings=[query_embedding],
        n_results=top_k
    )

    docs = results.get("documents", [[]])[0]
    metas = results.get("metadatas", [[]])[0]
    return list(zip(docs, metas))


# ---------------- OPENROUTER REQUEST ----------------
def make_openrouter_request(api_key: str, model: str, messages: List[Dict[str, str]]) -> str:
    headers = {
        "Authorization": f"Bearer {api_key}",
        "Content-Type": "application/json",
    }

    payload = {
        "model": model,
        "messages": messages,
        "temperature": 0.2,
    }

    response = requests.post(
        OPENROUTER_API_URL,
        headers=headers,
        json=payload,
        timeout=90,
    )

    response.raise_for_status()
    data = response.json()
    return data["choices"][0]["message"]["content"].strip()


# ---------------- BUILD PROMPT ----------------
def build_messages(chat_history: List[Dict[str, str]], question: str, contexts: List[str]) -> List[Dict[str, str]]:
    context_text = "\n\n".join(contexts) if contexts else "(No retrieved context)"

    system_prompt = (
        "You are a document QA assistant. Use only the provided context. "
        "If answer is not found in context, say you don't know."
    )

    messages: List[Dict[str, str]] = [
        {"role": "system", "content": system_prompt}
    ]

    messages.extend(chat_history[-8:])

    messages.append({
        "role": "user",
        "content": f"Context:\n{context_text}\n\nQuestion: {question}",
    })

    return messages


# ---------------- MAIN APP ----------------
def main():
    st.set_page_config(page_title="Advanced Local RAG", layout="wide")
    st.title("Advanced Local RAG (Streamlit + OpenRouter + ChromaDB)")

    embedder = get_embedder()
    collection = get_collection()

    default_api_key = os.getenv("OPENROUTER_API_KEY", "")
    default_model = os.getenv("OPENROUTER_MODEL", "stepfun/step-3.5-flash:free")

    if "messages" not in st.session_state:
        st.session_state.messages = []
    if "last_retrieved" not in st.session_state:
        st.session_state.last_retrieved = []

    with st.sidebar:
        st.subheader("LLM Settings")
        api_key = st.text_input("OpenRouter API Key", value=default_api_key, type="password")
        model_name = st.text_input("OpenRouter Model", value=default_model)
        top_k = st.slider("Retrieved chunks", 2, 8, 4)

        st.divider()
        st.subheader("Upload Documents")

        uploaded_files = st.file_uploader(
            "Upload files",
            accept_multiple_files=True,
            type=[
                "txt", "md", "pdf", "docx",
                "csv", "xlsx", "xls",
                "pptx", "json", "html", "zip"
            ],
        )

        if st.button("Index Files", type="primary"):
            if not uploaded_files:
                st.warning("Upload at least one file.")
            else:
                docs = []
                for file in uploaded_files:
                    text = read_uploaded_file(file)
                    docs.append((file.name, text))

                with st.spinner("Embedding and storing..."):
                    inserted = add_documents(collection, embedder, docs)

                st.success(f"Indexed {inserted} new chunks.")

        if st.button("Reset Vector DB"):
            client = chromadb.PersistentClient(
                path=DB_DIR,
                settings=Settings(anonymized_telemetry=False),
            )
            client.delete_collection(COLLECTION_NAME)
            client.get_or_create_collection(COLLECTION_NAME)
            st.success("Vector DB reset.")

        if st.button("Clear Chat"):
            st.session_state.messages = []
            st.session_state.last_retrieved = []
            st.success("Chat cleared.")

    # -------- CHAT --------
    st.subheader("Chat")

    for msg in st.session_state.messages:
        with st.chat_message(msg["role"]):
            st.write(msg["content"])

    question = st.chat_input("Ask about your documents...")

    if question:
        if not api_key.strip():
            st.warning("Add your OpenRouter API key.")
            st.stop()

        if collection.count() == 0:
            st.warning("No documents indexed yet.")
            st.stop()

        st.session_state.messages.append({"role": "user", "content": question})

        with st.chat_message("assistant"):
            with st.spinner("Retrieving context..."):
                retrieved = retrieve_context(collection, embedder, question, top_k)

            contexts = [doc for doc, _ in retrieved if doc]
            st.session_state.last_retrieved = retrieved

            llm_messages = build_messages(
                st.session_state.messages[:-1],
                question,
                contexts
            )

            with st.spinner("Generating answer..."):
                answer = make_openrouter_request(api_key, model_name, llm_messages)

            st.write(answer)

        st.session_state.messages.append({"role": "assistant", "content": answer})

    # -------- SHOW RETRIEVED --------
    if st.session_state.last_retrieved:
        st.markdown("### Retrieved Context")
        for idx, (doc, meta) in enumerate(st.session_state.last_retrieved, 1):
            source = meta.get("source", "unknown")
            chunk = meta.get("chunk", "?")
            with st.expander(f"{idx}. {source} (chunk {chunk})"):
                st.write(doc)


if __name__ == "__main__":
    os.makedirs(DB_DIR, exist_ok=True)
    main()